"""
PDF → Gemini(학과 추론 + 요약)
    → Gemma(질문 초안 생성)
    → Gemini(초안 관련성 검증 + 최소 수정)
    → [CLI 답변 → Gemini 평가 → Gemma 초안 → Gemini 수정] 루프

설치:
  pip install google-genai transformers peft torch accelerate pymupdf sentencepiece

사용법:
  python interview_question_generator.py --file 자료.pdf --adapter . --gemini-key YOUR_KEY
  python interview_question_generator.py --demo
"""

import argparse
import json
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path

from google import genai
from google.genai import types
from peft import PeftModel
from transformers import AutoTokenizer, AutoModelForCausalLM
import torch

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 설정
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

GEMINI_API_KEY    = os.environ.get("GEMINI_API_KEY", "YOUR_GEMINI_API_KEY_HERE")
LORA_ADAPTER_PATH = os.environ.get("LORA_ADAPTER_PATH", "./lora_adapter")
BASE_MODEL_ID     = "google/gemma-2b-it"

# 학습 데이터에서 추출한 system prompt — 원문 그대로
SYSTEM_PROMPT = (
    "당신은 대학 입시 면접관 AI입니다.\n"
    "Gemini 분석 결과와 직전 면접 맥락을 입력받아,\n"
    "다음 행동을 결정하고 JSON 패킷 하나만 출력하십시오.\n\n"
    "질문 방향 우선순위:\n"
    "  1순위: 지원자가 직접 수행한 탐구·실험·프로젝트·연구에 대한 질문\n"
    "         (예: 어떤 실험을 했는지, 결과가 예상과 달랐던 점, 한계와 개선 방향)\n"
    "  2순위: 탐구 과정에서 활용한 개념·원리에 대한 확인 질문\n"
    "         (원리 설명 자체가 목적이 아니라 탐구와 연결될 때만 사용)\n\n"
    "판단 기준:\n"
    "  follow_up  : 답변이 모호하거나 탐구 경험이 더 필요한 경우 → 날카로운 꼬리질문\n"
    "  next_topic : 답변이 충분히 구체적이거나 새 주제로 이동할 경우 → 자연스러운 전환\n\n"
    "출력은 반드시 유효한 JSON 하나만 생성하십시오. 설명/마크다운 절대 금지."
)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1단계: PDF 텍스트 추출
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def extract_pdf_text(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    return "\n".join(page.get_text() for page in doc)


def extract_pptx_text(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    lines = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"[슬라이드 {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_pdf_text(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx_text(file_path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext} (pdf, pptx만 가능)")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2단계: Gemini — PDF 초기 분석 (1회)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

INIT_PROMPT_TMPL = """당신은 대학 입시 면접 보조 AI입니다.
아래 문서를 분석하고, 반드시 순수 JSON만 출력하십시오.
마크다운/코드블록/설명 절대 금지. 출력의 첫 글자는 반드시 {{ 이어야 합니다.

필드 설명:
- dept: 지원 학과명 (문서에 없으면 내용 기반 추론)
- dept_reasoning: 학과 추론 근거 한 문장
- keywords: 면접 질문 소재가 될 구체적 기술 용어 10개 이상 (배열)
- topics: 면접에서 독립적으로 다룰 수 있는 주제 5개 이상. "주제명: 핵심 내용 한 문장" 형식 (배열)
- gemini_summary: 문서 전체 종합 요약. 각 개념의 핵심 원리·상호 관계·취약 포인트 포함. 600자 이내 한국어
- gemma_hint: Gemma 면접관 AI에게 전달할 첫 질문 생성용 힌트. 반드시 "'키워드' 관련 탐구·실험 경험 확인 필요. 직접 수행한 탐구가 있는지 질문." 형식으로 작성. 50자 이내.

출력 형식 예시:
{{"dept":"학과명","dept_reasoning":"근거","keywords":["키워드1","키워드2"],"topics":["주제1: 설명","주제2: 설명"],"gemini_summary":"요약","gemma_hint":"'합성곱' 키워드 언급했으나 메커니즘 설명 없음. 꼬리질문으로 검증 필요."}}

문서 내용:
{doc_text}
"""

def gemini_analyze_file(file_path: str) -> dict:
    """PDF/PPTX를 분석해 키워드·요약 반환 (면접 시작 시 1회 호출)"""
    client = genai.Client(api_key=GEMINI_API_KEY)

    try:
        raw_text = extract_text(file_path)
    except Exception as e:
        print(f"      [파일 추출 실패] {e}")
        raw_text = ""

    if not raw_text.strip():
        print(f"      [경고] 텍스트가 비어있습니다. 파일을 확인해주세요.")
        return {"dept": "미분류", "dept_reasoning": "분석 실패", "keywords": [], "gemini_summary": ""}

    prompt = INIT_PROMPT_TMPL.format(doc_text=raw_text[:12000])

    cfg = types.GenerateContentConfig(
        max_output_tokens=4096,
        temperature=0.2,
    )
    import time
    for attempt in range(3):
        try:
            print(f"      [Gemini] 분석 중... (시도 {attempt+1}/3)")
            response = client.models.generate_content(
                model="gemini-3.5-flash",
                contents=prompt,
                config=cfg,
            )
            if response.text is None:
                raise ValueError("response.text is None")
            return _parse_gemini_json(response.text)
        except Exception as e:
            err = str(e)
            if "429" in err and attempt < 2:
                wait = (attempt + 1) * 15
                print(f"      [할당량 초과] {wait}초 후 재시도...")
                time.sleep(wait)
            else:
                print(f"      [Gemini 호출 실패] {e}")
                return None
    return None


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3단계: Gemini — 답변 평가 (매 턴)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

EVAL_PROMPT_TMPL = """당신은 대학 입시 면접 보조 AI입니다.
지원자 답변을 평가하고, 반드시 순수 JSON만 출력하십시오.
마크다운/코드블록/설명 절대 금지. 출력의 첫 글자는 반드시 {{ 이어야 합니다.

[핵심 키워드]: {keywords}
[면접 주제 목록]:
{topics}
[문서 요약]: {doc_summary}
[면접관 질문]: {question}
[지원자 답변]: {answer}

판단 기준:
- follow_up: 답변이 모호하거나 깊이가 부족한 경우
- next_topic: 답변이 충분히 구체적인 경우 → 주제 목록에서 아직 안 다룬 주제 선택

출력 형식 예시:
{{"decision":"follow_up","gemini_summary":"'키워드' 키워드 언급했으나 설명 부족. 검증 필요. 또는 다음 주제: 주제명. 핵심 내용. 100자 이내."}}

gemini_summary는 반드시 Gemma 학습 형식에 맞게 작성하십시오:
- follow_up: "'[키워드]' 관련 탐구 경험 언급했으나 [구체성이 부족한 점]. 탐구 과정/결과 꼬리질문 필요."
- next_topic: "'[새 주제]' 관련 직접 수행한 탐구·실험이 있는지 확인 필요."
100자 이내로 작성하십시오.
"""

def gemini_evaluate_answer(
    keywords: list,
    topics: list,
    doc_summary: str,
    question: str,
    answer: str,
) -> dict:
    """지원자 답변을 평가해 follow_up/next_topic + 새 summary 반환"""
    client = genai.Client(api_key=GEMINI_API_KEY)

    prompt = EVAL_PROMPT_TMPL.format(
        keywords=", ".join(keywords),
        topics="\n".join(f"  - {t}" for t in topics),
        doc_summary=doc_summary,
        question=question,
        answer=answer,
    )

    import time
    cfg = types.GenerateContentConfig(
        max_output_tokens=1024,
        temperature=0.2,
    )
    for attempt in range(3):
        try:
            response = client.models.generate_content(
                model="gemini-3.5-flash",
                contents=prompt,
                config=cfg,
            )
            if response.text is None:
                raise ValueError("response.text is None")
            result = _parse_gemini_json(response.text)
            if result.get("decision") not in ("follow_up", "next_topic"):
                result["decision"] = "follow_up"
            return result
        except Exception as e:
            err = str(e)
            if "429" in err and attempt < 2:
                wait = (attempt + 1) * 15
                print(f"      [할당량 초과] {wait}초 후 재시도...")
                time.sleep(wait)
            else:
                print(f"      [Gemini 평가 실패] {e}")
                return {"decision": "follow_up", "gemini_summary": doc_summary}
    return {"decision": "follow_up", "gemini_summary": doc_summary}


def _parse_gemini_json(text: str) -> dict:
    """
    Gemini 응답에서 JSON 객체를 추출한다.
    gemini-3.5-flash 는 response_mime_type 설정에도
    'Here is the JSON:\n```json\n{...}\n```' 형태로 답할 수 있음.
    """
    text = text.strip()

    # 1) 직접 파싱
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # 2) 마크다운 코드블록 전체 제거 후 파싱
    cleaned = re.sub(r"```(?:json)?\s*", "", text)
    cleaned = re.sub(r"```", "", cleaned).strip()
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    # 3) 텍스트 안에서 { ... } 블록 직접 탐색 (첫 { 부터 마지막 } 까지)
    start = text.find("{")
    end   = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(text[start:end + 1])
        except json.JSONDecodeError:
            pass

    print(f"      [Gemini JSON 파싱 실패] 원본: {text[:200]}")
    return {}


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3-B단계: Gemini — Gemma 초안 질문 검증 및 최소 수정
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

REFINE_PROMPT_TMPL = """당신은 대학 입시 면접 질문 검수 AI입니다.
아래 지시를 따르고, 출력은 반드시 JSON 객체 하나만 작성하십시오.
{{ 로 시작해서 }} 로 끝나야 합니다.
"Here is", "```", 설명문, 마크다운 등 JSON 외 어떤 텍스트도 절대 출력하지 마십시오.

[PDF 핵심 키워드]: {keywords}
[면접 주제 목록]:
{topics}
[Gemma 초안 질문]: {draft_question}

검수 규칙:
0. 질문 텍스트 정제 (항상 적용, 해당 표현이 있으면 반드시 modified: true):
   - "(웃으면서)", "(미소지으며)", "(끄덕이며)" 등 괄호 안 행동 묘사 → 모두 제거
   - "시간이 다 돼서~", "그럼~", "음~", "자~" 등 불필요한 도입 문구 → 제거
   - "이거 말고", "또 다른 질문" 등 맥락 없는 메타 발언 → 제거
   - 정제 후 질문 핵심만 남기고 자연스럽게 다듬을 것
1. 정제 후 질문이 위 키워드/주제와 관련 있고 탐구·경험·연구를 묻는 방향이면 → modified: true (정제만 적용)
2. 정제 후에도 키워드/주제와 무관하거나 단순 원리 설명만 요구하면
   → 어투·길이·압박 강도는 유지하되, 탐구·실험·경험을 묻는 방향으로 추가 교정 (modified: true)
   - 예) "CNN의 합성곱 원리를 설명해보세요" → "CNN을 활용한 탐구나 실험을 직접 해본 적 있나요?"
   - 예) "LSTM 게이트 구조를 설명하세요" → "LSTM 관련 프로젝트나 실험에서 겪은 어려움이 있었나요?"
   - 예) "RNN의 한계를 말해보세요" → "RNN의 한계를 직접 탐구하거나 느꼈던 경험이 있다면 말해보세요."
3. 절대로 완전히 새로운 질문을 창작하지 마십시오. 초안을 기반으로만 수정하십시오.

출력 예시 (이 형식 그대로, 다른 텍스트 없이):
{{"modified": false, "question": "초안 질문 그대로"}}
{{"modified": true, "question": "수정된 질문"}}
"""

# fallback 질문 풀 — topics에서 매번 다른 항목 사용
_fallback_topic_idx = 0

def _make_fallback_question(keywords: list, topics: list) -> str:
    """topics를 순환하며 fallback 질문 생성 (매 호출마다 다른 주제)"""
    global _fallback_topic_idx
    if topics:
        topic = topics[_fallback_topic_idx % len(topics)]
        _fallback_topic_idx += 1
        # topics 형식: "주제명: 핵심 내용" — 주제명만 추출
        topic_name = topic.split(":")[0].strip()
        return f"{topic_name}에 대해 구체적으로 설명해보세요."
    if keywords:
        kw = keywords[_fallback_topic_idx % len(keywords)]
        _fallback_topic_idx += 1
        return f"{kw}의 핵심 원리에 대해 설명해보세요."
    return "방금 답변하신 내용을 좀 더 구체적으로 설명해주시겠어요?"


def gemini_refine_question(
    keywords: list,
    topics: list,
    draft_question: str,
) -> str:
    """Gemma 초안 질문을 PDF 키워드 기준으로 검증·최소 수정 후 최종 질문 반환"""
    if not draft_question or draft_question == "질문을 생성하지 못했습니다.":
        return _make_fallback_question(keywords, topics)

    client = genai.Client(api_key=GEMINI_API_KEY)

    topic_str = "\n".join(f"  - {t}" for t in topics[:5])
    kw_str    = ", ".join(keywords[:10])

    prompt = REFINE_PROMPT_TMPL.format(
        keywords=kw_str,
        topics=topic_str,
        draft_question=draft_question,
    )

    import time
    cfg = types.GenerateContentConfig(
        max_output_tokens=2048,
        temperature=0.1,
    )
    for attempt in range(3):
        try:
            response = client.models.generate_content(
                model="gemini-3.5-flash",
                contents=prompt,
                config=cfg,
            )
            if response.text is None:
                print(f"      [Gemini 검수 실패] response.text=None — fallback 사용")
                return _make_fallback_question(keywords, topics)

            raw    = response.text.strip()
            result = _parse_gemini_json(raw)

            if result:
                refined      = result.get("question", "").strip()
                was_modified = result.get("modified", False)

                if not refined:
                    return _make_fallback_question(keywords, topics)

                if was_modified:
                    print(f"      [Gemini 수정] 초안: \"{draft_question}\"")
                    print(f"      [Gemini 수정] 수정: \"{refined}\"")
                else:
                    print(f"      [Gemini 검수] 초안 통과 (수정 없음)")

                return refined

            # JSON 파싱 완전 실패
            print(f"      [Gemini 검수 파싱 실패] 원본 앞 100자: {raw[:100]!r}")
            fallback = _make_fallback_question(keywords, topics)
            print(f"      [Gemini 검수 실패] topics fallback 사용: {fallback!r}")
            return fallback

        except Exception as e:
            err = str(e)
            if "429" in err and attempt < 2:
                wait = (attempt + 1) * 15
                print(f"      [할당량 초과] {wait}초 후 재시도...")
                time.sleep(wait)
            else:
                print(f"      [Gemini 검수 실패] {e} — fallback 사용")
                return _make_fallback_question(keywords, topics)

    return _make_fallback_question(keywords, topics)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 4단계: Gemma LoRA 로드 (싱글턴)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

_model     = None
_tokenizer = None

def load_gemma_lora(adapter_path: str):
    global _model, _tokenizer
    if _model is not None:
        return _model, _tokenizer

    print(f"      베이스 모델 : {BASE_MODEL_ID}")
    print(f"      LoRA 어댑터 : {adapter_path}")

    dtype      = torch.float16 if torch.cuda.is_available() else torch.float32
    device_map = "auto"        if torch.cuda.is_available() else "cpu"

    tokenizer = AutoTokenizer.from_pretrained(BASE_MODEL_ID, use_fast=True, legacy=False)
    tokenizer.pad_token      = tokenizer.eos_token
    tokenizer.padding_side   = "left"
    tokenizer.add_special_tokens({"additional_special_tokens": ["<start_of_turn>", "<end_of_turn>"]})

    base  = AutoModelForCausalLM.from_pretrained(BASE_MODEL_ID, dtype=dtype, device_map=device_map)
    model = PeftModel.from_pretrained(base, adapter_path)
    model.eval()

    _model, _tokenizer = model, tokenizer
    print("      [모델 로드 완료]")
    return model, tokenizer


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 5단계: 프롬프트 빌드
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def build_prompt(
    prev_question: str,
    gemini_summary: str,
    keywords: list = None,
    topics: list = None,
) -> str:
    # 학습 데이터 형식에 맞게 gemini_summary를 150자로 압축
    summary_short = gemini_summary[:150].rstrip()

    # 키워드 최대 8개, 주제 최대 3개만 포함 (토큰 예산 절약)
    kw_str    = ", ".join((keywords or [])[:8])
    topic_str = "\n".join(f"  - {t}" for t in (topics or [])[:3])

    user_content = (
        f"핵심 키워드: {kw_str}\n"
        f"면접 주제:\n{topic_str}\n"
        f"직전 면접관 질문: {prev_question}\n"
        f"Gemini 분석 결과: {summary_short}"
    )
    return (
        f"<start_of_turn>system\n{SYSTEM_PROMPT}<end_of_turn>\n"
        f"<start_of_turn>user\n{user_content}<end_of_turn>\n"
        f"<start_of_turn>model\n"
    )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 6단계: Gemma 추론
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def run_gemma(prompt_text: str, adapter_path: str, max_new_tokens: int = 300) -> tuple[dict, str]:
    model, tokenizer = load_gemma_lora(adapter_path)

    inputs = tokenizer(
        prompt_text, return_tensors="pt", truncation=True, max_length=1024,
    ).to(model.device)

    eos_ids = [tokenizer.eos_token_id]
    eot_id  = tokenizer.convert_tokens_to_ids("<end_of_turn>")
    if eot_id and eot_id != tokenizer.eos_token_id:
        eos_ids.append(eot_id)

    with torch.no_grad():
        outputs = model.generate(
            **inputs,
            max_new_tokens=max_new_tokens,
            do_sample=True,
            temperature=0.7,
            top_p=0.9,
            repetition_penalty=1.1,
            pad_token_id=tokenizer.pad_token_id,
            eos_token_id=eos_ids,
        )

    generated = outputs[0][inputs["input_ids"].shape[1]:]
    raw_text  = tokenizer.decode(generated, skip_special_tokens=True).strip()

    cleaned    = re.sub(r"```(?:json)?|```", "", raw_text).strip()
    json_match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if json_match:
        try:
            return json.loads(json_match.group()), raw_text
        except json.JSONDecodeError:
            pass

    # JSON 파싱 실패 — 첫 줄만 질문 텍스트로 사용 (raw 전체 넘기지 않음)
    first_line = raw_text.split("\n")[0].strip()
    fallback_q = first_line if first_line else "답변 내용에 대해 더 구체적으로 설명해보세요."
    return {
        "type": "server_content", "message_id": str(uuid.uuid4()),
        "content": {
            "text":    fallback_q,
            "emotion": {"label": "중립/전환", "score": 0.5, "intensity": "medium"},
        },
    }, raw_text


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 7단계: 출력 패킷 조립
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _emotion_to_action(label: str) -> str:
    for key, action in {
        "날카로움": "avatar_stern",   "압박":   "avatar_stern",
        "정중함":   "avatar_nod",     "마무리":  "avatar_nod",
        "중립":     "avatar_neutral", "전환":   "avatar_neutral",
        "격려":     "avatar_smile",   "지지":   "avatar_smile",
        "당혹":     "avatar_tilt",    "재질문":  "avatar_tilt",
    }.items():
        if key in label:
            return action
    return "avatar_neutral"

def format_output(raw_dict: dict, gemini_info: dict, decision: str = "") -> dict:
    content = raw_dict.get("content", {})
    emotion = content.get("emotion", {})
    return {
        "type":       "server_content",
        "message_id": raw_dict.get("message_id", str(uuid.uuid4())),
        "content": {
            "text":     content.get("text", ""),
            "decision": decision or content.get("decision", "follow_up"),
            "emotion": {
                "label":     emotion.get("label", "중립/전환"),
                "score":     round(float(emotion.get("score", 0.7)), 2),
                "intensity": emotion.get("intensity", "medium"),
                "action":    _emotion_to_action(emotion.get("label", "")),
            },
        },
        "gemini_analysis": {
            "dept":     gemini_info.get("dept", ""),
            "keywords": gemini_info.get("keywords", []),
            "summary":  gemini_info.get("gemini_summary", ""),
        },
        "usage": {"timestamp": datetime.now(timezone.utc).isoformat()},
    }


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 메인 파이프라인 (인터랙티브 루프)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def run_pipeline(file_path: str, adapter_path: str = LORA_ADAPTER_PATH):
    print("\n" + "═" * 62)
    print("  Gemini → Gemma LoRA  면접 시뮬레이터")
    print("  종료: q 입력 후 Enter")
    print("═" * 62)

    # ── 초기 PDF 분석 (1회) ────────────────────────────────────────────────
    print(f"\n[초기 분석] 파일 읽는 중... ({Path(file_path).name})")
    gemini_info = gemini_analyze_file(file_path)
    if gemini_info is None:
        print("\n[오류] 분석 실패. API 키 할당량을 확인하거나 잠시 후 다시 시도하세요.")
        return

    dept        = gemini_info.get("dept", "미분류")
    keywords    = gemini_info.get("keywords", [])
    topics      = gemini_info.get("topics", [])
    doc_summary = gemini_info.get("gemini_summary", "")

    print(f"  ▸ 학과    : {dept}")
    print(f"  ▸ 키워드  : {', '.join(keywords)}")
    print(f"  ▸ 주제 목록:")
    for t in topics:
        print(f"      - {t}")
    print(f"  ▸ 요약    : {doc_summary}")

    # ── Gemma 모델 미리 로드 ──────────────────────────────────────────────
    print(f"\n[모델 로드] Gemma-2b-it + LoRA 로딩 중...")
    load_gemma_lora(adapter_path)

    # ── 첫 질문 생성 ─────────────────────────────────────────────────────
    print(f"\n[질문 생성] 첫 번째 질문 생성 중...")
    gemma_hint = gemini_info.get("gemma_hint", "")
    if not gemma_hint and keywords:
        gemma_hint = f"'{keywords[0]}' 키워드 언급했으나 메커니즘 설명 없음. 꼬리질문으로 검증 필요."

    prompt = build_prompt(
        prev_question="",
        gemini_summary=gemma_hint,
        keywords=keywords,
        topics=topics,
    )
    result_dict, raw = run_gemma(prompt, adapter_path)

    print(f"\n[질문 검수] Gemini 검수 중...")
    draft_q  = result_dict.get("content", {}).get("text", raw)
    final_q  = gemini_refine_question(keywords, topics, draft_q)
    result_dict.setdefault("content", {})["text"] = final_q

    packet      = format_output(result_dict, gemini_info, decision="follow_up")

    turn             = 1
    current_summary  = doc_summary
    current_question = packet["content"]["text"]

    # ── 인터랙티브 루프 ───────────────────────────────────────────────────
    while True:
        print("\n" + "─" * 62)
        print(f"  [Q{turn}] {current_question}")
        print("─" * 62)
        print(json.dumps(packet, ensure_ascii=False, indent=2))
        print()

        try:
            answer = input("▶ 답변 입력 (종료: q): ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n[면접 종료]")
            break

        if answer.lower() == "q":
            print("\n[면접 종료]")
            break

        if not answer:
            print("  (답변이 없습니다. 다시 입력해주세요.)")
            continue

        # ── Gemini: 답변 평가 ─────────────────────────────────────────────
        print(f"\n[Gemini 평가] 답변 분석 중...")
        eval_result = gemini_evaluate_answer(
            keywords=keywords,
            topics=topics,
            doc_summary=doc_summary,
            question=current_question,
            answer=answer,
        )
        decision        = eval_result.get("decision", "follow_up")
        current_summary = eval_result.get("gemini_summary", current_summary)

        label = "꼬리질문" if decision == "follow_up" else "새 주제"
        print(f"  ▸ 판단    : {decision}  ({label})")
        print(f"  ▸ 새 요약 : {current_summary}")

        # ── Gemma: 다음 질문 생성 ─────────────────────────────────────────
        turn += 1
        print(f"\n[질문 생성] Q{turn} 생성 중...")

        prompt = build_prompt(
            prev_question=current_question,
            gemini_summary=current_summary,
            keywords=keywords,
            topics=topics,
        )
        result_dict, raw = run_gemma(prompt, adapter_path)

        print(f"\n[질문 검수] Gemini 검수 중...")
        draft_q  = result_dict.get("content", {}).get("text", raw)
        final_q  = gemini_refine_question(keywords, topics, draft_q)
        result_dict.setdefault("content", {})["text"] = final_q

        packet           = format_output(result_dict, gemini_info, decision=decision)
        current_question = packet["content"]["text"]


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 데모 모드
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def demo_without_file():
    print("\n[DEMO MODE] 구조 확인용 — API/모델 호출 없음\n")
    mock_info = {
        "dept": "인공지능학과",
        "keywords": ["CNN", "RNN", "LSTM"],
        "gemini_summary": "CNN 합성곱 연산, RNN 순환 구조, LSTM 게이트 메커니즘 설명. LSTM Cell State 기울기 소실 해결 방식 검증 필요.",
    }
    mock_raw = {
        "type": "server_content", "message_id": str(uuid.uuid4()),
        "content": {
            "text":    "LSTM에서 Cell State가 기울기 소실 문제를 해결하는 원리를 수식 수준으로 설명해보세요.",
            "emotion": {"label": "날카로움/압박", "score": 0.88, "intensity": "high"},
        },
    }
    packet = format_output(mock_raw, mock_info, decision="follow_up")
    print("[Q1]", packet["content"]["text"])
    print(json.dumps(packet, ensure_ascii=False, indent=2))

    print("\n--- 가상 답변: 'forget gate가 이전 상태를 지웁니다' ---")
    mock_eval = {"decision": "follow_up", "gemini_summary": "forget gate 언급했으나 수식 설명 없음. 구체적 게이트 연산 검증 필요."}
    print(f"Gemini 판단: {mock_eval['decision']}")
    print(f"새 요약: {mock_eval['gemini_summary']}")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CLI
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PDF → Gemini → Gemma LoRA 면접 시뮬레이터")
    parser.add_argument("--file",       type=str, help="입력 파일 경로 (PDF 또는 PPTX)")
    parser.add_argument("--adapter",    type=str, default=LORA_ADAPTER_PATH, help="LoRA 어댑터 폴더")
    parser.add_argument("--gemini-key", type=str, default="", help="Gemini API 키")
    parser.add_argument("--demo",       action="store_true", help="데모 모드")
    args = parser.parse_args()

    if args.demo:
        demo_without_file()
    else:
        if args.gemini_key:
            GEMINI_API_KEY = args.gemini_key
        if not args.file:
            print("오류: --file 로 PDF 또는 PPTX 파일을 지정하세요. (또는 --demo)")
            parser.print_help()
            exit(1)
        run_pipeline(file_path=args.file, adapter_path=args.adapter)