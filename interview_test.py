"""
면접 질문 생성 테스트 파이프라인
PDF/PPTX 파일 → Gemini API (핵심 요약 + 질문 포인트) → Gemma-2b-it LoRA → JSON 출력

사용법:
  python interview_test.py --file 자소서.pdf --dept 컴퓨터공학과
  python interview_test.py --file 발표자료.pptx --dept AI융합학과 --prev "머신러닝을 공부했다고요?"
"""

import argparse
import base64
import json
import os
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path

# ── 의존성 안내 ────────────────────────────────────────────────────────────
# pip install google-generativeai transformers peft torch accelerate pymupdf python-pptx
# ────────────────────────────────────────────────────────────────────────────

import google.generativeai as genai
from peft import PeftModel
from transformers import AutoModelForCausalLM, AutoTokenizer
import torch

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 설정
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "YOUR_GEMINI_API_KEY_HERE")
LORA_ADAPTER_PATH = os.environ.get("LORA_ADAPTER_PATH", "./lora_adapter")  # adapter_config.json 폴더
BASE_MODEL_ID = "google/gemma-2b-it"

EMOTION_LABELS = ["날카로움/압박", "정중함/마무리", "중립/전환", "격려/지지", "당혹/재질문"]


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 1단계: 파일 → Gemini 분석
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def extract_text_from_file(file_path: str) -> tuple[str, str]:
    """PDF/PPTX/TXT 파일에서 텍스트 추출. (raw_text, mime_type) 반환"""
    p = Path(file_path)
    ext = p.suffix.lower()

    if ext == ".pdf":
        import fitz  # pymupdf
        doc = fitz.open(file_path)
        text = "\n".join(page.get_text() for page in doc)
        return text, "application/pdf"

    elif ext in (".pptx", ".ppt"):
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"[슬라이드 {i}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides.append("\n".join(parts))
        return "\n\n".join(slides), "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    elif ext in (".txt", ".md"):
        text = p.read_text(encoding="utf-8")
        return text, "text/plain"

    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}  (pdf / pptx / txt 만 가능)")


def gemini_analyze(file_path: str, dept: str) -> str:
    """
    Gemini로 파일을 분석하여 면접용 핵심 요약 + 질문 포인트를 생성.
    1024 토큰 이내로 압축된 'Gemini 분석 결과' 문자열 반환.
    """
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel("gemini-2.0-flash")

    raw_text, mime_type = extract_text_from_file(file_path)

    # 파일을 직접 업로드하는 방식 (PDF/PPTX는 Gemini File API 사용)
    p = Path(file_path)
    ext = p.suffix.lower()

    prompt = f"""당신은 대학 입시 면접 보조 AI입니다.
지원 학과: {dept}

아래 문서를 분석하여 면접관이 사용할 '분석 결과'를 **한국어로 200자 이내**로 작성하세요.
다음 형식을 반드시 따르세요:
- 핵심 키워드 2~3개 나열
- 모호하거나 심층 검증이 필요한 부분 1~2개 지적
- follow_up 또는 next_topic 중 권고 행동 명시
- 마크다운, 번호 목록 금지 / 단문으로 작성

문서 내용:
{raw_text[:6000]}
"""

    # 파일 업로드가 가능한 포맷이면 File API, 아니면 텍스트만 전송
    try:
        if ext in (".pdf", ".pptx"):
            uploaded = genai.upload_file(file_path, mime_type=mime_type)
            response = model.generate_content(
                [uploaded, prompt],
                generation_config=genai.GenerationConfig(
                    max_output_tokens=300,
                    temperature=0.4,
                ),
            )
        else:
            response = model.generate_content(
                prompt,
                generation_config=genai.GenerationConfig(
                    max_output_tokens=300,
                    temperature=0.4,
                ),
            )
    except Exception as e:
        print(f"[Gemini 파일 API 실패, 텍스트 모드로 폴백] {e}")
        response = model.generate_content(
            prompt,
            generation_config=genai.GenerationConfig(
                max_output_tokens=300,
                temperature=0.4,
            ),
        )

    return response.text.strip()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 2단계: Gemma LoRA 모델 로드 (싱글턴)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

_model = None
_tokenizer = None


def load_gemma_lora(adapter_path: str):
    global _model, _tokenizer

    if _model is not None:
        return _model, _tokenizer

    print(f"[모델 로드] 베이스 모델: {BASE_MODEL_ID}")
    print(f"[모델 로드] LoRA 어댑터: {adapter_path}")

    device_map = "auto" if torch.cuda.is_available() else "cpu"
    dtype = torch.float16 if torch.cuda.is_available() else torch.float32

    tokenizer = AutoTokenizer.from_pretrained(
        adapter_path,           # tokenizer_config.json 위치
        local_files_only=False,
    )

    base = AutoModelForCausalLM.from_pretrained(
        BASE_MODEL_ID,
        torch_dtype=dtype,
        device_map=device_map,
    )

    model = PeftModel.from_pretrained(base, adapter_path)
    model.eval()

    _model, _tokenizer = model, tokenizer
    print("[모델 로드 완료]")
    return model, tokenizer


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 3단계: 프롬프트 빌드 (학습 데이터와 동일한 chat template)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

SYSTEM_PROMPT = (
    "당신은 대학 입시 면접관 AI입니다.\n"
    "Gemini 음성 분석 결과와 직전 면접 맥락을 입력받아,\n"
    "다음 행동을 결정하고 JSON 패킷 하나만 출력하십시오.\n\n"
    "판단 기준:\n"
    "  follow_up  : 답변이 모호하거나 핵심 키워드 검증이 필요한 경우 → 날카로운 꼬리질문\n"
    "  next_topic : 답변이 충분히 구체적이거나 새 섹션으로 이동할 경우 → 자연스러운 전환\n\n"
    "출력은 반드시 유효한 JSON 하나만 생성하십시오. 설명/마크다운 절대 금지."
)


def build_prompt(dept: str, prev_question: str, gemini_analysis: str) -> str:
    user_content = (
        f"학과: {dept}\n"
        f"직전 면접관 질문: {prev_question}\n"
        f"Gemini 분석 결과: {gemini_analysis}"
    )
    return (
        f"<start_of_turn>system\n{SYSTEM_PROMPT}<end_of_turn>\n"
        f"<start_of_turn>user\n{user_content}<end_of_turn>\n"
        f"<start_of_turn>model\n"
    )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 4단계: Gemma 추론 + JSON 파싱 + 출력 포맷
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def run_gemma(prompt_text: str, adapter_path: str, max_new_tokens: int = 256) -> dict:
    model, tokenizer = load_gemma_lora(adapter_path)

    inputs = tokenizer(
        prompt_text,
        return_tensors="pt",
        truncation=True,
        max_length=1024,
    ).to(model.device)

    with torch.no_grad():
        outputs = model.generate(
            **inputs,
            max_new_tokens=max_new_tokens,
            do_sample=True,
            temperature=0.7,
            top_p=0.9,
            repetition_penalty=1.1,
            pad_token_id=tokenizer.eos_token_id,
            eos_token_id=tokenizer.convert_tokens_to_ids("<end_of_turn>"),
        )

    generated = outputs[0][inputs["input_ids"].shape[1]:]
    raw_text = tokenizer.decode(generated, skip_special_tokens=True).strip()

    # JSON 추출
    json_match = re.search(r"\{.*\}", raw_text, re.DOTALL)
    if json_match:
        try:
            return json.loads(json_match.group()), raw_text
        except json.JSONDecodeError:
            pass

    # 파싱 실패 시 fallback 구조 반환
    return {
        "type": "server_content",
        "message_id": str(uuid.uuid4()),
        "content": {
            "text": raw_text or "질문을 생성하지 못했습니다.",
            "decision": "follow_up",
            "emotion": {"label": "중립/전환", "score": 0.5, "intensity": "medium"},
            "is_final": False,
        },
        "stt_result": "[지원자 답변]",
        "usage": {"timestamp": datetime.now(timezone.utc).isoformat()},
    }, raw_text


def format_output(raw_dict: dict, gemini_summary: str) -> dict:
    """학습 데이터 형식에 맞게 출력 패킷 정규화"""
    content = raw_dict.get("content", {})
    emotion = content.get("emotion", {})

    return {
        "type": "server_content",
        "message_id": raw_dict.get("message_id", str(uuid.uuid4())),
        "content": {
            "text": content.get("text", ""),
            "decision": content.get("decision", "follow_up"),
            "emotion": {
                "label": emotion.get("label", "중립/전환"),
                "score": round(float(emotion.get("score", 0.7)), 2),
                "intensity": emotion.get("intensity", "medium"),
                "action": _emotion_to_action(emotion.get("label", "")),
            },
            "is_final": content.get("is_final", False),
        },
        "gemini_analysis": gemini_summary,
        "stt_result": "[지원자 답변]",
        "usage": {
            "timestamp": datetime.now(timezone.utc).isoformat(),
        },
    }


def _emotion_to_action(label: str) -> str:
    mapping = {
        "날카로움": "avatar_stern",
        "압박": "avatar_stern",
        "정중함": "avatar_nod",
        "마무리": "avatar_nod",
        "중립": "avatar_neutral",
        "전환": "avatar_neutral",
        "격려": "avatar_smile",
        "지지": "avatar_smile",
        "당혹": "avatar_tilt",
        "재질문": "avatar_tilt",
    }
    for key, action in mapping.items():
        if key in label:
            return action
    return "avatar_neutral"


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 메인 파이프라인
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def run_pipeline(
    file_path: str,
    dept: str,
    prev_question: str = "",
    adapter_path: str = LORA_ADAPTER_PATH,
    verbose: bool = True,
) -> dict:
    """전체 파이프라인 실행 후 최종 JSON 반환"""

    print("\n" + "═" * 60)
    print("  면접 질문 생성 파이프라인")
    print("═" * 60)

    # Step 1: Gemini 분석
    print(f"\n[1/3] Gemini 파일 분석 중... ({Path(file_path).name})")
    gemini_result = gemini_analyze(file_path, dept)
    if verbose:
        print(f"      Gemini 분석 결과:\n      {gemini_result}\n")

    # Step 2: 프롬프트 빌드
    print("[2/3] Gemma 프롬프트 빌드 중...")
    prompt = build_prompt(dept, prev_question, gemini_result)
    if verbose:
        print("      프롬프트 (일부):")
        print("      " + prompt[:300].replace("\n", "\n      "))

    # Step 3: Gemma 추론
    print("\n[3/3] Gemma-2b-it (LoRA) 추론 중...")
    result_dict, raw_output = run_gemma(prompt, adapter_path)
    if verbose:
        print(f"      Raw 출력: {raw_output[:200]}")

    # 출력 포맷 정규화
    final_output = format_output(result_dict, gemini_result)

    print("\n" + "─" * 60)
    print("  최종 출력 패킷")
    print("─" * 60)
    print(json.dumps(final_output, ensure_ascii=False, indent=2))
    print("═" * 60 + "\n")

    return final_output


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 데모 모드 (실제 파일 없이 구조 확인용)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def demo_without_file():
    """파일·모델 없이 파이프라인 구조만 확인하는 데모"""
    print("\n[DEMO MODE] 실제 API/모델 호출 없이 구조 확인")

    mock_gemini = (
        "'딥러닝', '역전파', 'CNN' 키워드 언급. "
        "역전파 메커니즘 설명이 표면적 수준에 머뭄. "
        "꼬리질문으로 수식 수준 이해도 검증 필요. follow_up 권고."
    )
    mock_raw = json.dumps({
        "type": "server_content",
        "message_id": str(uuid.uuid4()),
        "content": {
            "text": "역전파에서 기울기 소실 문제를 어떻게 해결할 수 있나요?",
            "decision": "follow_up",
            "emotion": {"label": "날카로움/압박", "score": 0.85, "intensity": "high"},
            "is_final": False,
        },
        "stt_result": "[지원자 답변]",
        "usage": {"timestamp": datetime.now(timezone.utc).isoformat()},
    }, ensure_ascii=False)

    parsed, _ = json.loads(mock_raw), mock_raw
    final = format_output(parsed, mock_gemini)

    print("\n[Gemini 분석 결과]")
    print(f"  {mock_gemini}\n")
    print("[최종 출력 패킷]")
    print(json.dumps(final, ensure_ascii=False, indent=2))
    return final


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CLI
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Gemini + Gemma-2b-it LoRA 면접 질문 생성기")
    parser.add_argument("--file", type=str, help="입력 파일 경로 (PDF/PPTX/TXT)")
    parser.add_argument("--dept", type=str, default="컴퓨터공학과", help="지원 학과")
    parser.add_argument("--prev", type=str, default="", help="직전 면접관 질문")
    parser.add_argument("--adapter", type=str, default=LORA_ADAPTER_PATH, help="LoRA 어댑터 폴더 경로")
    parser.add_argument("--gemini-key", type=str, default="", help="Gemini API 키")
    parser.add_argument("--demo", action="store_true", help="데모 모드 (파일/모델 없이 구조 확인)")
    args = parser.parse_args()

    if args.demo:
        demo_without_file()
    else:
        if args.gemini_key:
            GEMINI_API_KEY = args.gemini_key

        if not args.file:
            print("오류: --file 옵션으로 입력 파일을 지정하세요. (또는 --demo로 구조만 확인)")
            parser.print_help()
            exit(1)

        run_pipeline(
            file_path=args.file,
            dept=args.dept,
            prev_question=args.prev,
            adapter_path=args.adapter,
        )